import os
import json  # <-- IMPORT THIS
import ast
import datetime
from typing import List, Dict, Tuple, Optional, Type
import re

from pinecone import Pinecone
from sentence_transformers import SentenceTransformer, CrossEncoder
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain_classic.agents import AgentExecutor, Tool, create_react_agent
from langchain_classic import hub
from langchain.tools import BaseTool
from pydantic import BaseModel, Field
from typing import Type
from reportlab.lib.pagesizes import letter
import markdown
from html2docx import html2docx
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import smtplib
from email.message import EmailMessage

# ... (Input Schemas like MDToolInput, DocxToolInput, etc. are unchanged) ...
class MDToolInput(BaseModel):
    """Input schema for the SaveAsMDTool."""
    content: str = Field(description="The text content to be written to the Markdown file.")
    filename: str = Field(description="The desired name for the output Markdown file, e.g., 'notes.md'.")

class DocxToolInput(BaseModel):
    """Input schema for the SaveAsDocxTool."""
    content: str = Field(description="The text content to be written to the DOCX file.")
    filename: str = Field(description="The desired name for the output DOCX file, e.g., 'report.docx'.")

class PptxToolInput(BaseModel):
    """Input schema for the SaveAsPptxTool."""
    content: str = Field(description="The text content to be written to the PPTX file. Newlines (\\n\\n) will be treated as new slides.")
    filename: str = Field(description="The desired name for the output PPTX file, e.g., 'presentation.pptx'.")

class PDFToolInput(BaseModel):
    """Input schema for the SaveAsPDFTool."""
    content: str = Field(description="The text content to be written to the PDF.")
    filename: str = Field(description="The desired name for the output PDF file, e.g., 'assignment.pdf'.")

class FeedbackInput(BaseModel):
    """Input schema for the FeedbackTool."""
    rating: str = Field(description="The user's rating, e.g., 'good' or 'bad'.")
    comment: Optional[str] = Field(
        default=None, 
        description="The user's free-text comment."
    )
    new_source: Optional[str] = Field(
        default=None, 
        description="A URL or path to a new information source suggested by the user."
    )
    
class EmailToolInput(BaseModel):
    """Input schema for the EmailTool."""
    to_email: str = Field(description="The recipient's email address.")
    subject: str = Field(description="The subject line of the email.")
    body: str = Field(description="The main text content of the email.")
    attachment_path: Optional[str] = Field(
        default=None, 
        description="The local file path to a document to attach (e.g., 'generated_documents/assignment.pdf')."
    )

# --- Tool Implementations with Download Fix ---

class SaveAsMDTool(BaseTool):
    """A tool that saves given text content into a Markdown (.md) file."""
    name: str = "Save As Markdown"
    description: str = "Use this tool to save text content into a .md file. It takes content and a filename as input."
    args_schema: Type[BaseModel] = MDToolInput

    def _run(self, content: str, filename: str) -> str:
        """Saves the content to a Markdown file and returns the status."""
        try:
            output_dir = "generated_documents"
            os.makedirs(output_dir, exist_ok=True)
            if not filename.endswith(".md"):
                filename += ".md"
            file_path = os.path.join(output_dir, filename)

            with open(file_path, "w", encoding="utf-8") as f:
                f.write(content)

            # --- *** START OF DOWNLOAD FIX *** ---
            # Return a JSON string with the download path
            download_path = f"{file_path}"
            response_data = {
                "download_path": download_path,
                "message": f"Successfully saved content to {filename}."
            }
            return json.dumps(response_data)
            # --- *** END OF DOWNLOAD FIX *** ---
            
        except Exception as e:
            return f"Error saving Markdown file: {e}"


class SaveAsPDFTool(BaseTool):
    """A tool that saves given text content into a PDF file."""
    name: str = "Save As PDF"
    description: str = "Use this tool to save text content into a .pdf file. It takes content and a filename as input."
    args_schema: Type[BaseModel] = PDFToolInput

    def _run(self, content: str, filename: str) -> str:
        """Saves the content to a PDF and returns the status."""
        try:
            output_dir = "generated_documents"
            os.makedirs(output_dir, exist_ok=True)
            if not filename.endswith(".pdf"):  # <-- Fix: ensure .pdf
                filename += ".pdf"
            file_path = os.path.join(output_dir, filename)
            
            html_content = markdown.markdown(content)
            doc = SimpleDocTemplate(file_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            for line in html_content.splitlines():
                story.append(Paragraph(line, styles["Normal"]))
                story.append(Spacer(1, 12))
            doc.build(story)
            
            # --- *** START OF DOWNLOAD FIX *** ---
            download_path = f"{file_path}"
            response_data = {
                "download_path": download_path,
                "message": f"Successfully saved content to {filename}."
            }
            return json.dumps(response_data)
            # --- *** END OF DOWNLOAD FIX *** ---
            
        except Exception as e:
            return f"Error saving PDF: {e}"

class SaveAsDocxTool(BaseTool):
    """A tool that saves given text content into a Word (.docx) file."""
    name: str = "Save As DOCX"
    description: str = "Use this tool to save text content into a .docx file. It takes content and a filename as input."
    args_schema: Type[BaseModel] = DocxToolInput

    def _run(self, content: str, filename: str) -> str:
        """Saves the content to a DOCX file and returns the status."""
        try:
            output_dir = "generated_documents"
            os.makedirs(output_dir, exist_ok=True)
            if not filename.endswith(".docx"):
                filename += ".docx"
            file_path = os.path.join(output_dir, filename)
            
            doc = Document()
            doc.add_heading(filename.replace(".docx", ""), 0)
            
            for para in content.split('\n\n'):
                doc.add_paragraph(para)
                
            doc.save(file_path)

            # --- *** START OF DOWNLOAD FIX *** ---
            download_path = f"{file_path}"
            response_data = {
                "download_path": download_path,
                "message": f"Successfully saved content to {filename}."
            }
            return json.dumps(response_data)
            # --- *** END OF DOWNLOAD FIX *** ---
            
        except Exception as e:
            return f"Error saving DOCX file: {e}"

class SaveAsPptxTool(BaseTool):
    """A tool that saves given text content into a PowerPoint (.pptx) file."""
    name: str = "Save As PPTX"
    description: str = "Use this tool to save text content into a .pptx presentation file. It takes content and a filename as input."
    args_schema: Type[BaseModel] = PptxToolInput

    def _run(self, content: str, filename: str) -> str:
        """Saves the content to a PPTX file, splitting slides by Markdown headings."""
        try:
            output_dir = "generated_documents"
            os.makedirs(output_dir, exist_ok=True)
            if not filename.endswith(".pptx"):
                filename += ".pptx"
            file_path = os.path.join(output_dir, filename)
            
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            content_slide_layout = prs.slide_layouts[1]
            
            slide_chunks = re.split(r'\n(?=#{1,6} )', content.strip())

            if not slide_chunks or (len(slide_chunks) == 1 and not slide_chunks[0].strip()):
                return "Error: Content was empty or had no headings to split into slides."

            for i, slide_md in enumerate(slide_chunks):
                # ... (rest of the PPTX generation logic is unchanged) ...
                slide_md = slide_md.strip()
                if not slide_md:
                    continue
                
                lines = slide_md.split('\n', 1)
                title_line = lines[0].strip()
                body_content = lines[1].strip() if len(lines) > 1 else ""
                
                slide_title = title_line.lstrip('# ').strip()
                
                if i == 0:
                    slide = prs.slides.add_slide(title_slide_layout)
                    slide.shapes.title.text = slide_title
                    if body_content:
                        try:
                            slide.placeholders[1].text = body_content
                        except:
                            pass 
                
                else:
                    slide = prs.slides.add_slide(content_slide_layout)
                    slide.shapes.title.text = slide_title
                    
                    if body_content:
                        content_shape = slide.shapes.placeholders[1]
                        tf = content_shape.text_frame
                        tf.clear() 

                        current_p = None
                        for line in body_content.split('\n'):
                            line_stripped = line.strip()
                            if not line_stripped:
                                continue

                            indentation = len(line) - len(line.lstrip(' '))
                            level = indentation // 4 

                            if line_stripped.startswith(('* ', '- ', '1. ')):
                                text = re.sub(r'^(\* |-\ |\d+\. )', '', line_stripped)
                                
                                p = tf.add_paragraph()
                                p.text = text
                                p.level = level
                            else:
                                p = tf.add_paragraph()
                                p.text = line_stripped
                                p.level = level

            if len(prs.slides) == 0:
                 return "Error: No slides were generated. The content might be missing markdown headings."
                
            prs.save(file_path)

            # --- *** START OF DOWNLOAD FIX *** ---
            download_path = f"{file_path}"
            response_data = {
                "download_path": download_path,
                "message": f"Successfully saved content to {filename}."
            }
            return json.dumps(response_data)
            # --- *** END OF DOWNLOAD FIX *** ---
            
        except Exception as e:
            return f"Error saving PPTX file: {e}"

# ... (FeedbackTool and EmailTool are unchanged) ...
class FeedbackTool(BaseTool):
    """A tool that records user feedback for later review."""
    name: str = "Record Feedback"
    description: str = "Use this tool to record user feedback (good or bad), including comments or new information sources."
    args_schema: Type[BaseModel] = FeedbackInput
    
    def _run(self, rating: str, comment: Optional[str] = None, new_source: Optional[str] = None) -> str:
        """Logs the feedback to a file."""
        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(script_dir)
            log_file = os.path.join(project_root, "feedback_log.jsonl")

            feedback_entry = {
                "timestamp": datetime.datetime.now().isoformat(),
                "rating": rating,
                "comment": comment, 
                "new_source": new_source
            }
            with open(log_file, "a") as f:
                f.write(json.dumps(feedback_entry) + "\n")
            
            return "Thank you, your feedback has been recorded."
        except Exception as e:
            return f"Error recording feedback: {e}"
            
class EmailTool(BaseTool):
    """A tool that sends an email, optionally with an attachment."""
    name: str = "Send Email"
    description: str = "Use this tool to send an email. It can optionally attach one file from the 'generated_documents' folder."
    args_schema: Type[BaseModel] = EmailToolInput

    def _run(self, to_email: str, subject: str, body: str, attachment_path: Optional[str] = None) -> str:
        """Sends the email and returns the status."""
        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(script_dir)
            assets_path = os.path.join(project_root, "assets.json")
            
            with open(assets_path, 'r') as f:
                credentials = json.load(f)
            
            SENDER_EMAIL = credentials.get("SENDER_EMAIL")
            SENDER_EMAIL_PASSWORD = credentials.get("SENDER_EMAIL_PASSWORD")

            if not SENDER_EMAIL or not SENDER_EMAIL_PASSWORD:
                return "Error: SENDER_EMAIL or SENDER_EMAIL_PASSWORD not found in assets.json"

            msg = EmailMessage()
            msg['Subject'] = subject
            msg['From'] = SENDER_EMAIL
            msg['To'] = to_email
            msg.set_content(body)

            if attachment_path:
                if not os.path.exists(attachment_path):
                    return f"Error: Attachment path not found: {attachment_path}"
                
                import mimetypes
                ctype, encoding = mimetypes.guess_type(attachment_path)
                if ctype is None or encoding is not None:
                    ctype = 'application/octet-stream'
                
                maintype, subtype = ctype.split('/', 1)
                
                with open(attachment_path, 'rb') as fp:
                    msg.add_attachment(fp.read(),
                                       maintype=maintype,
                                       subtype=subtype,
                                       filename=os.path.basename(attachment_path))
                print(f"Attached file: {attachment_path}")

            print(f"Connecting to SMTP server to send email to {to_email}...")
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
                smtp.login(SENDER_EMAIL, SENDER_EMAIL_PASSWORD)
                smtp.send_message(msg)
            
            return f"Successfully sent email to {to_email} with subject: {subject}"

        except Exception as e:
            return f"Error sending email: {e}"
            
# ... (Rest of SMEAgent class is unchanged) ...
class SMEAgent:
    """
    A Subject Matter Expert agent that uses an advanced
    Parent-Child retrieval pipeline with reranking.
    """
    def __init__(self, assets_path: str = "assets.json"):
        print("Initializing SMEAgent...")
        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            self.parent_dir = os.path.dirname(script_dir) 
            self._load_credentials(os.path.join(self.parent_dir, assets_path))
        except NameError:
            self.parent_dir = os.path.abspath(os.path.join(os.getcwd(), os.pardir))
            self._load_credentials(os.path.join(self.parent_dir, assets_path))
            
        self.llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=self.GEMINI_API_KEY, temperature=0)
        
        print("Loading embedding and reranker models...")
        self.embeddings_model = SentenceTransformer('all-mpnet-base-v2')
        self.reranker = CrossEncoder('BAAI/bge-reranker-base', max_length=512)
        print("Models loaded.")

        self.pc = Pinecone(api_key=os.environ["PINECONE_API_KEY"])
        self.child_index = self.pc.Index("sme-agent-child-chunks")
        self.parent_index = self.pc.Index("sme-agent-parent-chunks")
        print("Pinecone indexes connected.")
        
        print("SMEAgent initialized successfully.")

    def _load_credentials(self, file_path: str):
        self.GEMINI_API_KEY = self._get_api_key("gemini_api_key", file_path)
        os.environ["PINECONE_API_KEY"] = self._get_api_key("pinecone_api_key", file_path)

        if not all([os.environ["PINECONE_API_KEY"], self.GEMINI_API_KEY]):
            raise ValueError("All API keys must be set in the assets file.")

    def _get_api_key(self, key_name: str, file_path: str):
        try:
            with open(file_path, 'r') as f:
                credentials = json.load(f)
            return credentials.get(key_name)
        except (FileNotFoundError, json.JSONDecodeError) as e:
            print(f"Error reading credentials file at {file_path}: {e}")
            return None

    def _create_tools(self) -> list:
        # This method isn't used by the graph, but is left for completeness
        return [
            Tool(
                name="Question Answering",
                func=self._run_qa_rag, 
                description="Use this tool ONLY to answer specific questions about data science concepts. Input should be a clear question."
            ),
            # ... other tools ...
            FeedbackTool()
        ]

    def _retrieve_and_rerank(self, query: str, final_top_k: int = 3) -> str:
        print(f"RAG: Retrieving context for query: '{query}'")
        query_embedding = self.embeddings_model.encode(query).tolist()
        retrieval_top_k = 20  
        child_results = self.child_index.query(
            vector=query_embedding,
            top_k=retrieval_top_k,
            include_metadata=True
        )
        parent_ids = []
        for match in child_results.get('matches', []):
            if 'parent_chunk_id' in match.get('metadata', {}):
                parent_ids.append(match['metadata']['parent_chunk_id'])
        parent_ids = list(set(parent_ids))
        if not parent_ids:
            print("RAG: No parent documents found from retrieval.")
            return "No relevant context found."
        print(f"RAG: Fetching {len(parent_ids)} unique parent documents...")
        parent_results = self.parent_index.fetch(ids=parent_ids)        
        docs_to_process = []
        for parent_id, doc_data in parent_results['vectors'].items():
            docs_to_process.append({
                "id": parent_id,
                "text": doc_data['metadata'].get('text', ''),
                "source": doc_data['metadata'].get('source', 'Unknown')
            })
        try:
            print(f"RAG: Reranking {len(docs_to_process)} candidate documents...")
            pairs = [[query, doc['text']] for doc in docs_to_process]
            scores = self.reranker.predict(pairs, show_progress_bar=False)            
            reranked_docs = list(zip(scores, docs_to_process))
            reranked_docs.sort(key=lambda x: x[0], reverse=True)            
            context_chunks = []
            print(f"RAG: Building context from top {final_top_k} reranked results.")
            for score, doc in reranked_docs[:final_top_k]:
                context_chunks.append(doc['text'])
            
            return "\n\n---\n\n".join(context_chunks)

        except Exception as e:
            print(f"RAG: Reranking failed ({e}). Falling back to basic similarity.")
            fallback_context = [doc['text'] for doc in docs_to_process[:final_top_k]]
            return "\n\n---\n\n".join(fallback_context)

    def _get_qa_prompt(self):
        template = """You are a helpful Data Science assistant.
                        Read the user's QUESTION and the provided CONTEXT carefully.
                        Think step-by-step to find the relevant information in the context that answers the question.

                        Finally, answer the QUESTION based ONLY on the CONTEXT and answer step-by-step.
                        If the context does not contain the answer, state that you cannot answer.

                        CONTEXT:
                        {context}

                        QUESTION:
                        {question}

                        ANSWER:"""
        return ChatPromptTemplate.from_template(template)

    def _get_assignment_prompt(self):
        template = """You are a Data Science instructor. Create the text content for an assignment based on the user's request, using the provided context.
                        Generate only the text for the assignment.

                        CONTEXT:
                        {context}

                        USER'S REQUEST:
                        {question}

                        ASSIGNMENT CONTENT:"""
        return ChatPromptTemplate.from_template(template)

    def run_qa_rag(self, query: str) -> str:
        print("--- Running QA RAG Chain ---")
        qa_prompt = self._get_qa_prompt()
        
        rag_chain = (
            {"context": self._retrieve_and_rerank, "question": RunnablePassthrough()}
            | qa_prompt
            | self.llm
            | StrOutputParser()
        )
        
        return rag_chain.invoke(query)
    
    def run_assignment_rag(self, query: str) -> str:
        print("--- Running Assignment RAG Chain ---")
        assignment_prompt = self._get_assignment_prompt()
        
        rag_chain = (
            {"context": self._retrieve_and_rerank, "question": RunnablePassthrough()}
            | assignment_prompt
            | self.llm
            | StrOutputParser()
        )
        
        return rag_chain.invoke(query)

    def run(self, user_input: str) -> dict:
        print(f"Agent is processing input: '{user_input}'")
        # This is part of the old agent executor, not used by the graph
        # self.agent_executor.invoke({"input": user_input})
        return {"output": "This method is not used by the graph."}