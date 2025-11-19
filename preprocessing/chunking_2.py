import os
import re
from typing import List, Dict, Tuple, Optional
import json
import time

# File loading libraries
import pypdf
import docx
from pptx import Presentation

# Pinecone
from pinecone import Pinecone, ServerlessSpec
from tqdm import tqdm

# Embeddings & Reranking
from sentence_transformers import SentenceTransformer, CrossEncoder # Added CrossEncoder

# LangChain document object
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter

def load_document(file_path: str) -> str:
    """
    Loads text content from a file based on its extension.
    Supports .pdf, .docx, .pptx, .txt, .md, and .py files.
    """
    print(f"Loading document: {file_path}")
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(file_path)
            text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text = "\n".join(text_runs)
        elif ext in [".txt", ".md", ".py"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            print(f"Unsupported file type: {ext}. Skipping.")
            return ""
        return text
    except Exception as e:
        print(f"Error loading {file_path}: {e}")
        return ""

def clean_text(text: str) -> str:
    """
    Cleans the extracted text with a more sophisticated approach.
    """
    
    # 1. Normalize ligatures (common in PDFs)
    text = text.replace("ﬁ", "fi").replace("ﬂ", "fl")

    # 2. De-hyphenate words broken across newlines
    text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
    
    # 3. Normalize whitespace characters
    text = text.replace('\t', ' ').replace('\xa0', ' ')

    text = re.sub(r' +', ' ', text)

    # 5. Handle newlines
    # 5a. Collapse 3 or more newlines into a double newline (paragraph break)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 5b. Remove single newlines that are not part of a paragraph break
    text = re.sub(r'([^\s])\n([^\s])', r'\1 \2', text)

    # 6. Clean up: Remove leading/trailing spaces on each line
    text = re.sub(r'[ \t]*\n[ \t]*', '\n', text)
    
    # 7. Strip leading/trailing whitespace from the whole text
    text = text.strip()
    
    return text

def chunk_text_parent_child(
    text: str, 
    file_path: str, 
    parent_chunk_size: int = 2000,
    parent_chunk_overlap: int = 200,
    child_chunk_size: int = 500,
    child_chunk_overlap: int = 50
) -> Tuple[List[str], List[str], Dict[int, int]]:
    """
    Splits text into parent and child chunks with a mapping between them.
    (Kept for compatibility, but main script now uses a more direct method)
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    parent_splitter = RecursiveCharacterTextSplitter(
        chunk_size=parent_chunk_size, 
        chunk_overlap=parent_chunk_overlap, 
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    child_splitter = RecursiveCharacterTextSplitter(
        chunk_size=child_chunk_size, 
        chunk_overlap=child_chunk_overlap, 
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    parent_docs = []
    if file_ext == ".md":
        headers_to_split_on = [("#", "Header 1"), ("##", "Header 2"), ("###", "Header 3")]
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
        md_header_splits = markdown_splitter.split_text(text)
        parent_docs = parent_splitter.split_documents(md_header_splits)
    else:
        parent_docs = [Document(page_content=chunk) for chunk in parent_splitter.split_text(text)]
    
    parent_chunks = [doc.page_content for doc in parent_docs]

    child_chunks = []
    child_to_parent_map = {}
    
    for parent_idx, parent_doc in enumerate(parent_docs):
        child_doc_chunks = child_splitter.split_documents([parent_doc])
        
        for child_doc in child_doc_chunks:
            child_idx = len(child_chunks)
            child_chunks.append(child_doc.page_content)
            child_to_parent_map[child_idx] = parent_idx
    
    print(f"Created {len(parent_chunks)} parent chunks and {len(child_chunks)} child chunks")
    
    return parent_chunks, child_chunks, child_to_parent_map


def get_pinecone_credentials_from_json(file_path: str):
    """Reads the Pinecone API key from a JSON credentials file."""
    try:
        # Get path relative to this script
        script_dir = os.path.dirname(os.path.abspath(__file__))
        full_path = os.path.join(script_dir, file_path)
        
        with open(full_path, 'r') as f:
            credentials = json.load(f)
        return credentials.get("pinecone_api_key")
    except Exception as e:
        print(f"Error loading Pinecone key from {full_path}: {e}")
        return None

# ================== MAIN SCRIPT ==================

if __name__ == "__main__":
    DATA_DIRECTORY = "../data"
    
    if not os.path.exists(DATA_DIRECTORY):
        print(f"Error: Directory '{DATA_DIRECTORY}' not found.")
        print("Please make sure the path is correct and you are running this from the correct directory.")
    else:

        PINECONE_API_KEY = get_pinecone_credentials_from_json("../assets.json")
        if not PINECONE_API_KEY:
            raise ValueError("Could not find 'pinecone_api_key' in ../assets.json. Please check the file.")
            
        pc = Pinecone(api_key=PINECONE_API_KEY)
        
        print("Initializing embeddings model (all-mpnet-base-v2)...")
        embedding_model_name = "all-mpnet-base-v2"
        embeddings_model = SentenceTransformer(embedding_model_name)
        EMBEDDING_DIM = embeddings_model.get_sentence_embedding_dimension()
        print(f"Embeddings Model loaded. Embedding dimension: {EMBEDDING_DIM}")

        print("Initializing BGE Reranker model (BAAI/bge-reranker-base)...")
        rerank_model = CrossEncoder('BAAI/bge-reranker-base', max_length=512)
        print("Reranker model loaded.")

        CHILD_INDEX_NAME = "sme-agent-child-chunks"
        PARENT_INDEX_NAME = "sme-agent-parent-chunks"
        
        # Create child index (for vector search)
        if CHILD_INDEX_NAME not in pc.list_indexes().names():
            print(f"Creating child index: {CHILD_INDEX_NAME}")
            pc.create_index(
                name=CHILD_INDEX_NAME,
                dimension=EMBEDDING_DIM,
                metric='cosine',
                spec=ServerlessSpec(cloud='aws', region='us-east-1')
            )
            print("Waiting 5s for child index to initialize...")
            time.sleep(5) 
        else:
             print(f"Child index '{CHILD_INDEX_NAME}' already exists.")
        
        # Create parent index (for key-value retrieval)
        if PARENT_INDEX_NAME not in pc.list_indexes().names():
            print(f"Creating parent index: {PARENT_INDEX_NAME}")
            pc.create_index(
                name=PARENT_INDEX_NAME,
                dimension=EMBEDDING_DIM, # Pinecone requires vectors, even for fetching
                metric='cosine',
                spec=ServerlessSpec(cloud='aws', region='us-east-1')
            )
            print("Waiting 5s for parent index to initialize...")
            time.sleep(5)
        else:
            print(f"Parent index '{PARENT_INDEX_NAME}' already exists.")

        
        child_index = pc.Index(CHILD_INDEX_NAME)
        parent_index = pc.Index(PARENT_INDEX_NAME)
        
        # --- 3. Define splitters ---
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000, 
            chunk_overlap=200, 
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        child_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500, 
            chunk_overlap=50, 
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        all_parent_docs_for_upload = []
        all_child_docs_for_upload = []
        doc_counter = 0
        print(f"Walking directory: {DATA_DIRECTORY}...")        
        for root, _, files in os.walk(DATA_DIRECTORY):
            for file in files:
                file_path = os.path.join(root, file)
                
                raw_text = load_document(file_path)
                if not raw_text:
                    continue
                
                cleaned_text = clean_text(raw_text)
                if not cleaned_text:
                    print(f"Skipping {file_path} after cleaning (empty).")
                    continue

                # Create parent LangChain Documents
                file_ext = os.path.splitext(file_path)[1].lower()
                current_file_parents = []
                
                if file_ext == ".md":
                    headers_to_split_on = [("#", "Header 1"), ("##", "Header 2"), ("###", "Header 3")]
                    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                    md_header_splits = markdown_splitter.split_text(cleaned_text)
                    current_file_parents = parent_splitter.split_documents(md_header_splits)
                else:
                    current_file_parents = [Document(page_content=chunk) for chunk in parent_splitter.split_text(cleaned_text)]
                
                for parent_doc in current_file_parents:
                    parent_id = f"doc_{doc_counter}"
                    parent_doc.metadata = {
                        "source": file_path,
                        "file_name": os.path.basename(file_path),
                        "parent_chunk_id": parent_id
                    }
                    all_parent_docs_for_upload.append(parent_doc)
                    
                    # Create child docs from this parent
                    child_doc_chunks = child_splitter.split_documents([parent_doc])
                    for i, child_chunk in enumerate(child_doc_chunks):
                        child_chunk.metadata = {
                            "source": file_path,
                            "file_name": os.path.basename(file_path),
                            "parent_chunk_id": parent_id,
                            "child_id": f"{parent_id}_child_{i}"
                        }
                        all_child_docs_for_upload.append(child_chunk)
                    
                    doc_counter += 1
                
                print(f"Processed {file}: {len(current_file_parents)} parent chunks")
        
        print(f"\nTotal parent documents to upload: {len(all_parent_docs_for_upload)}")
        print(f"Total child documents to upload: {len(all_child_docs_for_upload)}")
        
        if not all_parent_docs_for_upload:
            print("No documents processed. Proceeding to test query.")
        else:
            BATCH_SIZE = 100 # Batch size for upserting
            
            # --- 5. Upsert Child Chunks to Pinecone ---
            print(f"\nUpserting {len(all_child_docs_for_upload)} child documents to {CHILD_INDEX_NAME}...")
            for i in tqdm(range(0, len(all_child_docs_for_upload), BATCH_SIZE)):
                batch_docs = all_child_docs_for_upload[i:i + BATCH_SIZE]
                
                texts = [doc.page_content for doc in batch_docs]
                embeddings = embeddings_model.encode(texts).tolist()
                
                vectors_to_upsert = []
                for j, doc in enumerate(batch_docs):
                    pinecone_metadata = {
                        "source_file": doc.metadata.get("file_name", "unknown"),
                        "parent_chunk_id": doc.metadata["parent_chunk_id"]
                    }
                    vectors_to_upsert.append({
                        "id": doc.metadata["child_id"],
                        "values": embeddings[j],
                        "metadata": pinecone_metadata 
                    })
                child_index.upsert(vectors=vectors_to_upsert)

            # --- 6. Upsert Parent Chunks to Pinecone ---
            print(f"\nUpserting {len(all_parent_docs_for_upload)} parent documents to {PARENT_INDEX_NAME}...")
            for i in tqdm(range(0, len(all_parent_docs_for_upload), BATCH_SIZE)):
                batch_docs = all_parent_docs_for_upload[i:i + BATCH_SIZE]
                
                texts = [doc.page_content for doc in batch_docs]
                embeddings = embeddings_model.encode(texts).tolist() 
                
                vectors_to_upsert = []
                for j, doc in enumerate(batch_docs):
                    vectors_to_upsert.append({
                        "id": doc.metadata["parent_chunk_id"],
                        "values": embeddings[j],
                        "metadata": {
                            "text": doc.page_content,
                            "source": doc.metadata.get("file_name", "unknown")
                        }
                    })
                parent_index.upsert(vectors=vectors_to_upsert)
            
            print("\nAll documents upserted to Pinecone.")
            print("Waiting 10s for index to update...")
            time.sleep(10)
        
        # --- 7. Test the retriever with RERANKING and FALLBACK ---
        print("\n" + "=" * 80)
        print("Testing the Pinecone retriever with Reranking...")
        print("=" * 80)
        
        query = "What are loss functions?"
        final_top_k = 3             # How many final results to show
        retrieval_top_k = 10        # How many candidates to fetch for reranking
        
        print(f"Querying: '{query}'")
        print("-" * 80)
        
        # 1. Embed query
        query_embedding = embeddings_model.encode(query).tolist()
        
        # 2. Query child index
        print(f"Fetching top {retrieval_top_k} candidates from child index...")
        child_results = child_index.query(
            vector=query_embedding,
            top_k=retrieval_top_k,
            include_metadata=True
        )
        
        # 3. Get unique parent IDs
        parent_ids = []
        for match in child_results.get('matches', []):
            if 'parent_chunk_id' in match.get('metadata', {}):
                parent_ids.append(match['metadata']['parent_chunk_id'])
        parent_ids = list(set(parent_ids))
        
        if not parent_ids:
            print("No results found from initial retrieval.")
        else:
            # 4. Fetch parent chunks by ID
            print(f"Fetching {len(parent_ids)} unique parent documents...")
            parent_results = parent_index.fetch(ids=parent_ids)
            
            # Store parent docs for reranking or fallback
            docs_to_process = []
            for parent_id, doc_data in parent_results['vectors'].items():
                docs_to_process.append({
                    "id": parent_id,
                    "text": doc_data['metadata'].get('text', ''),
                    "source": doc_data['metadata'].get('source', 'Unknown')
                })

            # --- 5. RERANKING (try...except for fallback) ---
            try:
                print(f"Attempting to rerank {len(docs_to_process)} candidate documents...")
                # Create [query, passage] pairs
                pairs = [[query, doc['text']] for doc in docs_to_process]
                
                # Predict scores using the local CrossEncoder model
                scores = rerank_model.predict(pairs, show_progress_bar=True)
                
                # Combine scores with documents and sort
                reranked_docs = list(zip(scores, docs_to_process))
                reranked_docs.sort(key=lambda x: x[0], reverse=True)

                print(f"\nRetrieved Top {final_top_k} RERANKED results:")
                print("-" * 80)
                for i, (score, doc) in enumerate(reranked_docs[:final_top_k], 1):
                    print(f"\n[Reranked Document {i} (Score: {score:.4f})]")
                    print(f"Source: {doc['source']}")
                    print(f"Parent Chunk ID: {doc['id']}")
                    print(f"\nContent preview:\n{doc['text'][:400]}...")
                    print("-" * 80)

            except Exception as e:
                # --- FALLBACK LOGIC ---
                print(f"\n!!! Reranking failed ({e}). Falling back to basic similarity. !!!")
                print(f"\nRetrieved Top {final_top_k} (Basic Similarity) results:")
                print("-" * 80)
                # Use the parent docs as fetched (they are not sorted by relevance here)
                # For a better fallback, you could sort `docs_to_process` based on `child_results` scores
                # before slicing, but this is the simplest fallback.
                for i, doc in enumerate(docs_to_process[:final_top_k], 1):
                    print(f"\n[Fallback Document {i}]")
                    print(f"Source: {doc['source']}")
                    print(f"Parent Chunk ID: {doc['id']}")
                    print(f"\nContent preview:\n{doc['text'][:400]}...")
                    print("-" * 80)

        print("\n--- Final Index Statistics ---")
        try:
            print("\nChild index statistics:")
            print(child_index.describe_index_stats())
            print("\nParent index statistics:")
            print(parent_index.describe_index_stats())
        except Exception as e:
            print(f"Could not fetch index stats: {e}")