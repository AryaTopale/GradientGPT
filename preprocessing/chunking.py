import os
import re
import pickle
from typing import List, Dict, Tuple, Optional

# File loading libraries
import pypdf
import docx
from pptx import Presentation

# LangChain document object
from langchain_core.documents import Document
# In-memory store for parent documents
from langchain_classic.storage import InMemoryStore
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_classic.retrievers import ParentDocumentRetriever


def load_document(file_path: str) -> str:
    """
    Loads text content from a file based on its extension.
    Supports .pdf, .docx, .pptx, .txt, .md, and .py files.
    """
    print(f"Loading document: {file_path}")
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(file_path)
            text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text = "\n".join(text_runs)
        elif ext in [".txt", ".md", ".py"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            print(f"Unsupported file type: {ext}. Skipping.")
            return ""
        return text
    except Exception as e:
        print(f"Error loading {file_path}: {e}")
        return ""


def clean_text(text: str) -> str:
    """
    Cleans the extracted text with a more sophisticated approach.
    """
    # 1. Normalize ligatures (common in PDFs)
    text = text.replace("ﬁ", "fi").replace("ﬂ", "fl")

    # 2. De-hyphenate words broken across newlines
    text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
    
    # 3. Normalize whitespace characters
    text = text.replace('\t', ' ').replace('\xa0', ' ')
    
    # 4. Collapse multiple spaces into a single space
    text = re.sub(r' +', ' ', text)

    # 5. Handle newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'([^\s])\n([^\s])', r'\1 \2', text)

    # 6. Clean up: Remove leading/trailing spaces on each line
    text = re.sub(r'[ \t]*\n[ \t]*', '\n', text)
    
    # 7. Strip leading/trailing whitespace from the whole text
    text = text.strip()
    
    return text


def save_retriever(retriever: ParentDocumentRetriever, save_dir: str = "./saved_retriever"):
    """
    Saves the ParentDocumentRetriever to disk.
    
    Args:
        retriever: The ParentDocumentRetriever to save
        save_dir: Directory to save the retriever components
    """
    os.makedirs(save_dir, exist_ok=True)
    
    print(f"Saving retriever to {save_dir}...")
    
    # Save the FAISS vectorstore
    vectorstore_path = os.path.join(save_dir, "vectorstore")
    retriever.vectorstore.save_local(vectorstore_path)
    print(f"  ✓ Saved vectorstore")
    
    # Save the docstore (contains parent documents)
    docstore_path = os.path.join(save_dir, "docstore.pkl")
    with open(docstore_path, 'wb') as f:
        pickle.dump(retriever.docstore.store, f)
    print(f"  ✓ Saved docstore")
    
    # Save metadata about splitters
    metadata = {
        'child_chunk_size': retriever.child_splitter._chunk_size,
        'child_chunk_overlap': retriever.child_splitter._chunk_overlap,
        'parent_chunk_size': retriever.parent_splitter._chunk_size if hasattr(retriever.parent_splitter, '_chunk_size') else 2000,
        'parent_chunk_overlap': retriever.parent_splitter._chunk_overlap if hasattr(retriever.parent_splitter, '_chunk_overlap') else 200,
    }
    metadata_path = os.path.join(save_dir, "metadata.pkl")
    with open(metadata_path, 'wb') as f:
        pickle.dump(metadata, f)
    print(f"  ✓ Saved metadata")
    
    print("Retriever saved successfully!")


def load_retriever(
    save_dir: str = "./saved_retriever",
    embedding_model: str = "all-mpnet-base-v2",
    hf_token: Optional[str] = None
) -> ParentDocumentRetriever:
    """
    Loads a saved ParentDocumentRetriever from disk.
    
    Args:
        save_dir: Directory containing the saved retriever
        embedding_model: Name of the HuggingFace embedding model (must match the one used during saving)
        hf_token: Optional HuggingFace API token
        
    Returns:
        ParentDocumentRetriever ready for querying
    """
    print(f"Loading retriever from {save_dir}...")
    
    # Load metadata
    metadata_path = os.path.join(save_dir, "metadata.pkl")
    with open(metadata_path, 'rb') as f:
        metadata = pickle.load(f)
    print(f"  ✓ Loaded metadata")
    
    # Recreate embeddings
    embeddings = HuggingFaceEmbeddings(
        model_name=embedding_model,
        model_kwargs={'token': hf_token} if hf_token else {}
    )
    
    # Load the FAISS vectorstore
    vectorstore_path = os.path.join(save_dir, "vectorstore")
    vectorstore = FAISS.load_local(
        vectorstore_path, 
        embeddings,
        allow_dangerous_deserialization=True
    )
    print(f"  ✓ Loaded vectorstore")
    
    # Load the docstore
    docstore_path = os.path.join(save_dir, "docstore.pkl")
    docstore = InMemoryStore()
    with open(docstore_path, 'rb') as f:
        docstore.store = pickle.load(f)
    print(f"  ✓ Loaded docstore")
    
    # Recreate splitters
    child_splitter = RecursiveCharacterTextSplitter(
        chunk_size=metadata['child_chunk_size'],
        chunk_overlap=metadata['child_chunk_overlap'],
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    parent_splitter = RecursiveCharacterTextSplitter(
        chunk_size=metadata['parent_chunk_size'],
        chunk_overlap=metadata['parent_chunk_overlap'],
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    # Create the retriever
    retriever = ParentDocumentRetriever(
        vectorstore=vectorstore,
        docstore=docstore,
        child_splitter=child_splitter,
        parent_splitter=parent_splitter,
    )
    
    print("Retriever loaded successfully!")
    return retriever


def create_parent_child_retriever(
    text: str,
    file_path: str,
    parent_chunk_size: int = 2000,
    parent_chunk_overlap: int = 200,
    child_chunk_size: int = 500,
    child_chunk_overlap: int = 50,
    embedding_model: str = "all-mpnet-base-v2",
    hf_token: Optional[str] = None
) -> ParentDocumentRetriever:
    """
    Creates a ParentDocumentRetriever from text.
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # Create parent documents
    if file_ext == ".md":
        headers_to_split_on = [
            ("#", "Header 1"),
            ("##", "Header 2"),
            ("###", "Header 3"),
        ]
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
        md_header_splits = markdown_splitter.split_text(text)
        
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=parent_chunk_size,
            chunk_overlap=parent_chunk_overlap,
        )
        parent_docs = parent_splitter.split_documents(md_header_splits)
    else:
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=parent_chunk_size,
            chunk_overlap=parent_chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        parent_docs = [Document(page_content=chunk) for chunk in parent_splitter.split_text(text)]
    
    # Add metadata to parent documents
    for i, doc in enumerate(parent_docs):
        doc.metadata.update({
            "source": file_path,
            "file_name": os.path.basename(file_path),
        })
    
    # Setup splitters
    child_splitter = RecursiveCharacterTextSplitter(
        chunk_size=child_chunk_size,
        chunk_overlap=child_chunk_overlap
    )
    
    # Setup embeddings
    embeddings = HuggingFaceEmbeddings(
        model_name=embedding_model,
        model_kwargs={'token': hf_token} if hf_token else {}
    )
    
    # Create empty vectorstore
    vectorstore = FAISS.from_texts(["dummy"], embeddings)
    
    # Create in-memory docstore
    docstore = InMemoryStore()
    
    # Create retriever
    retriever = ParentDocumentRetriever(
        vectorstore=vectorstore,
        docstore=docstore,
        child_splitter=child_splitter,
        parent_splitter=parent_splitter,
    )
    
    # Add parent documents
    retriever.add_documents(parent_docs, ids=None)
    
    print(f"Created retriever with {len(parent_docs)} parent documents")
    
    return retriever


def query_parent_child_retriever(
    retriever: ParentDocumentRetriever, 
    query: str, 
    k: int = 3
) -> List[Document]:
    """
    Queries the parent-child retriever and returns parent documents.
    """
    results = retriever.invoke(query, k=k)
    return results


if __name__ == "__main__":
    DATA_DIRECTORY = "./data"
    SAVED_RETRIEVER_DIR = "./saved_retriever"
    
    # Check if saved retriever exists
    if os.path.exists(SAVED_RETRIEVER_DIR):
        print("=" * 80)
        print("Found saved retriever. Loading from disk...")
        print("=" * 80)
        retriever = load_retriever(SAVED_RETRIEVER_DIR)
        
    else:
        print("=" * 80)
        print("No saved retriever found. Creating new one...")
        print("=" * 80)
        
        if not os.path.exists(DATA_DIRECTORY):
            print(f"Error: Directory '{DATA_DIRECTORY}' not found.")
            exit(1)
        
        # Define splitters
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000, 
            chunk_overlap=200, 
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        child_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500, 
            chunk_overlap=50, 
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        all_parent_docs = []
        
        print(f"Walking directory: {DATA_DIRECTORY}...")
        
        # Walk directory, load, clean, and create parent docs
        for root, _, files in os.walk(DATA_DIRECTORY):
            for file in files:
                file_path = os.path.join(root, file)
                
                raw_text = load_document(file_path)
                if not raw_text:
                    continue
                
                cleaned_text = clean_text(raw_text)
                if not cleaned_text:
                    print(f"Skipping {file_path} after cleaning (empty).")
                    continue

                # Split text into parent documents
                file_ext = os.path.splitext(file_path)[1].lower()
                current_file_parents = []
                
                if file_ext == ".md":
                    headers_to_split_on = [("#", "Header 1"), ("##", "Header 2"), ("###", "Header 3")]
                    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                    md_header_splits = markdown_splitter.split_text(cleaned_text)
                    current_file_parents = parent_splitter.split_documents(md_header_splits)
                else:
                    current_file_parents = [Document(page_content=chunk) for chunk in parent_splitter.split_text(cleaned_text)]
                
                # Add metadata
                for doc in current_file_parents:
                    doc.metadata.update({
                        "source": file_path,
                        "file_name": os.path.basename(file_path)
                    })
                    all_parent_docs.append(doc)
                
                print(f"Processed {file}: {len(current_file_parents)} parent chunks")
        
        print(f"\nTotal parent documents created: {len(all_parent_docs)}")
        
        if not all_parent_docs:
            print("No documents processed. Exiting.")
            exit(1)
        
        # Initialize embeddings and vectorstore
        print("Initializing embeddings model (all-mpnet-base-v2)...")
        embeddings = HuggingFaceEmbeddings(model_name="all-mpnet-base-v2")
        
        print("Creating empty FAISS vectorstore...")
        vectorstore = FAISS.from_texts(["dummy"], embeddings)
        
        print("Creating in-memory docstore...")
        docstore = InMemoryStore()

        # Create the retriever
        retriever = ParentDocumentRetriever(
            vectorstore=vectorstore,
            docstore=docstore,
            child_splitter=child_splitter,
            parent_splitter=parent_splitter,
        )
        
        print("Adding documents to retriever...")
        retriever.add_documents(all_parent_docs, ids=None)
        
        print("\nRetriever created successfully!")
        
        # Save the retriever
        save_retriever(retriever, SAVED_RETRIEVER_DIR)
    
    # Test the retriever
    print("\n" + "=" * 80)
    print("Testing the retriever with a query...")
    print("=" * 80)
    
    query = "What are loss functions?"
    print(f"Querying: '{query}'")
    print("-" * 80)
    
    results = query_parent_child_retriever(retriever, query, k=3)
    
    if not results:
        print("No results found. The data might not contain info on this query.")
    else:
        for i, doc in enumerate(results, 1):
            print(f"\n[Parent Document {i}]")
            print(f"Source: {doc.metadata.get('file_name', 'Unknown')}")
            print(f"Parent Chunk ID: {doc.metadata.get('doc_id', 'Unknown')}")
            print(f"\nContent preview:\n{doc.page_content[:400]}...")
            print("-" * 80)